from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
BLACK   = RGBColor(0x0a, 0x0a, 0x0a)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIME    = RGBColor(0xAA, 0xFF, 0x00)
LIME_DK = RGBColor(0x66, 0x99, 0x00)
GREY_1  = RGBColor(0x1a, 0x1a, 0x1a)
GREY_2  = RGBColor(0x2a, 0x2a, 0x2a)
GREY_3  = RGBColor(0x44, 0x44, 0x44)
DIM     = RGBColor(0x88, 0x88, 0x88)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_name='Georgia', size=24, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, l, t, w, h,
                  font_name='Calibri', size=16, color=WHITE,
                  align=PP_ALIGN.LEFT, line_spacing=1.2):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb

def slide_bg(slide, color=BLACK):
    bg = add_rect(slide, 0, 0, W, H, fill=color)
    return bg

def add_lime_rule(slide, l, t, w=Inches(0.4), h=Pt(3)):
    add_rect(slide, l, t, w, Emu(int(h)), fill=LIME)

def add_pill(slide, text, l, t, font_size=9):
    pw = Inches(1.6)
    ph = Inches(0.28)
    add_rect(slide, l, t, pw, ph, fill=GREY_2)
    add_text(slide, text, l, t, pw, ph,
             font_name='Courier New', size=font_size, color=LIME,
             align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — The Opportunity
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
slide_bg(s1)

# Subtle grid texture (thin lines)
for i in range(1, 8):
    add_rect(s1, Inches(i * 13.33/8), 0, Emu(6000), H, fill=RGBColor(0x18,0x18,0x18))
for j in range(1, 5):
    add_rect(s1, 0, Inches(j * 7.5/5), W, Emu(6000), fill=RGBColor(0x18,0x18,0x18))

# Left accent stripe
add_rect(s1, 0, 0, Inches(0.06), H, fill=LIME)

# Slide number
add_text(s1, '01 / 04', Inches(0.3), Inches(0.3), Inches(1.2), Inches(0.3),
         font_name='Courier New', size=8, color=DIM)

# Eyebrow
add_pill(s1, 'THE OPPORTUNITY', Inches(0.3), Inches(1.1))

# Main headline
add_text(s1, 'Great characters\ndon\'t scale\nthemselves.',
         Inches(0.3), Inches(1.55), Inches(5.8), Inches(2.8),
         font_name='Georgia', size=52, bold=False, color=WHITE)

# Lime rule
add_lime_rule(s1, Inches(0.3), Inches(4.55), Inches(0.5))

# Body copy
body = [
    'IAMS Proactive Health has a genuinely original brand character — the Gut Bacteria.',
    '',
    'But turning that character into a constant stream of social content means',
    'briefing studios, waiting weeks, and paying for every single execution.',
    '',
    'The volume social demands today makes that model impossible to sustain.',
]
add_multiline(s1, body, Inches(0.3), Inches(4.7), Inches(5.6), Inches(2.3),
              font_name='Calibri', size=14, color=DIM)

# Right visual — stat cards
for i, (num, label) in enumerate([
        ('6–8 wks', 'avg. studio turnaround'),
        ('$12–40K', 'cost per video execution'),
        ('3–5×/wk', 'content cadence required'),
    ]):
    cy = Inches(1.4 + i * 1.85)
    add_rect(s1, Inches(7.2), cy, Inches(5.5), Inches(1.55), fill=GREY_1)
    add_rect(s1, Inches(7.2), cy, Inches(0.06), Inches(1.55), fill=LIME)
    add_text(s1, num, Inches(7.4), cy + Inches(0.15), Inches(5.0), Inches(0.7),
             font_name='Georgia', size=38, bold=False, color=WHITE)
    add_text(s1, label.upper(), Inches(7.4), cy + Inches(0.75), Inches(5.0), Inches(0.5),
             font_name='Courier New', size=9, color=DIM)

# Footer
add_rect(s1, 0, Inches(7.18), W, Emu(40000), fill=GREY_2)
add_text(s1, 'IAMS PROACTIVE HEALTH  ·  CONFIDENTIAL',
         Inches(0.3), Inches(7.2), Inches(6), Inches(0.28),
         font_name='Courier New', size=7, color=GREY_3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Solution
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
slide_bg(s2)

add_rect(s2, 0, 0, Inches(0.06), H, fill=LIME)
add_text(s2, '02 / 04', Inches(0.3), Inches(0.3), Inches(1.2), Inches(0.3),
         font_name='Courier New', size=8, color=DIM)

add_pill(s2, 'THE SOLUTION', Inches(0.3), Inches(1.1))

add_text(s2, 'The Insiders\nFactory.',
         Inches(0.3), Inches(1.55), Inches(6.2), Inches(2.2),
         font_name='Georgia', size=56, bold=False, color=WHITE)

add_lime_rule(s2, Inches(0.3), Inches(3.9), Inches(0.5))

add_text(s2, 'A purpose-built content studio for IAMS — in a browser tab.',
         Inches(0.3), Inches(4.1), Inches(5.8), Inches(0.5),
         font_name='Calibri', size=16, color=WHITE)

add_multiline(s2, [
    'Any brand or agency team member can generate a finished,',
    'voice-matched, lip-synced character video in under ten minutes.',
    '',
    'No production brief. No studio. No waiting.',
    'Just a direction, a voice line, and a format.',
], Inches(0.3), Inches(4.7), Inches(5.6), Inches(2.2),
    font_name='Calibri', size=13, color=DIM)

# Right — feature tiles
features = [
    ('On-brand by design',    'The character, voice, and visual brief are locked in.\nEvery output is brand-safe before it renders.'),
    ('Voice + lip-sync',      'The character speaks your chosen line.\nMouth movement is frame-accurate, every time.'),
    ('Any format, instantly', '9:16 for Stories. 1:1 for Feed. 16:9 for YouTube.\nSwitch with a single click.'),
    ('Zero technical skill',  'Choose a mood preset or write a custom direction.\nGenerate. Download. Post.'),
]
for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    fx = Inches(7.0 + col * 3.15)
    fy = Inches(1.1 + row * 2.9)
    add_rect(s2, fx, fy, Inches(2.9), Inches(2.55), fill=GREY_1)
    add_rect(s2, fx, fy, Inches(2.9), Emu(18000), fill=GREY_2)
    add_text(s2, title, fx + Inches(0.18), fy + Inches(0.12), Inches(2.55), Inches(0.35),
             font_name='Calibri', size=11, bold=True, color=LIME)
    add_multiline(s2, desc.split('\n'), fx + Inches(0.18), fy + Inches(0.52),
                  Inches(2.55), Inches(1.85), font_name='Calibri', size=11, color=DIM)

add_rect(s2, 0, Inches(7.18), W, Emu(40000), fill=GREY_2)
add_text(s2, 'IAMS PROACTIVE HEALTH  ·  CONFIDENTIAL',
         Inches(0.3), Inches(7.2), Inches(6), Inches(0.28),
         font_name='Courier New', size=7, color=GREY_3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — How It Works
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
slide_bg(s3)

add_rect(s3, 0, 0, Inches(0.06), H, fill=LIME)
add_text(s3, '03 / 04', Inches(0.3), Inches(0.3), Inches(1.2), Inches(0.3),
         font_name='Courier New', size=8, color=DIM)

add_pill(s3, 'HOW IT WORKS', Inches(0.3), Inches(1.1))

add_text(s3, 'Brief to broadcast\nin minutes.',
         Inches(0.3), Inches(1.55), Inches(6.5), Inches(1.8),
         font_name='Georgia', size=48, color=WHITE)

add_lime_rule(s3, Inches(0.3), Inches(3.5), Inches(0.5))

add_text(s3, 'A four-step automated pipeline — no human handoffs required.',
         Inches(0.3), Inches(3.72), Inches(6.0), Inches(0.45),
         font_name='Calibri', size=14, color=DIM)

# Pipeline steps — horizontal across the bottom half
steps = [
    ('01', 'DIRECTION',  'Choose a mood preset\nor write a custom\ndirection for the scene.'),
    ('02', 'VOICE LINE', 'Select a brand-approved\nscript line. AI generates\na matched voice track.'),
    ('03', 'VIDEO',      'The character is rendered\nin motion — bouncy,\nexpressive, on-brand.'),
    ('04', 'LIP-SYNC',   'Audio is fused to the\nvideo frame-by-frame.\nReady to download.'),
]
step_w = Inches(2.8)
step_h = Inches(2.8)
step_gap = Inches(0.35)
step_start_x = Inches(0.3)
step_y = Inches(4.1)

for i, (num, title, desc) in enumerate(steps):
    sx = step_start_x + i * (step_w + step_gap)
    # Card bg
    add_rect(s3, sx, step_y, step_w, step_h, fill=GREY_1)
    # Top lime stripe
    add_rect(s3, sx, step_y, step_w, Emu(14000), fill=LIME if i == 0 else GREY_2)
    # Step number
    add_text(s3, num, sx + Inches(0.15), step_y + Inches(0.1), Inches(0.5), Inches(0.4),
             font_name='Courier New', size=10,
             color=BLACK if i == 0 else LIME)
    # Arrow connector (not last)
    if i < 3:
        ax = sx + step_w + Inches(0.05)
        add_text(s3, '→', ax, step_y + step_h/2 - Inches(0.25), Inches(0.28), Inches(0.5),
                 font_name='Courier New', size=18, color=GREY_3, align=PP_ALIGN.CENTER)
    # Title
    add_text(s3, title, sx + Inches(0.15), step_y + Inches(0.52), step_w - Inches(0.3), Inches(0.45),
             font_name='Courier New', size=10, bold=True,
             color=BLACK if i == 0 else LIME)
    # Desc
    add_multiline(s3, desc.split('\n'), sx + Inches(0.15), step_y + Inches(1.05),
                  step_w - Inches(0.3), Inches(1.5),
                  font_name='Calibri', size=12, color=DIM if i > 0 else RGBColor(0x22,0x22,0x22))

add_rect(s3, 0, Inches(7.18), W, Emu(40000), fill=GREY_2)
add_text(s3, 'IAMS PROACTIVE HEALTH  ·  CONFIDENTIAL',
         Inches(0.3), Inches(7.2), Inches(6), Inches(0.28),
         font_name='Courier New', size=7, color=GREY_3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Impact
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
slide_bg(s4)

add_rect(s4, 0, 0, Inches(0.06), H, fill=LIME)
add_text(s4, '04 / 04', Inches(0.3), Inches(0.3), Inches(1.2), Inches(0.3),
         font_name='Courier New', size=8, color=DIM)

add_pill(s4, 'THE IMPACT', Inches(0.3), Inches(1.1))

add_text(s4, 'Always-on content.\nAlways brand-safe.',
         Inches(0.3), Inches(1.55), Inches(6.2), Inches(2.0),
         font_name='Georgia', size=48, color=WHITE)

add_lime_rule(s4, Inches(0.3), Inches(3.65), Inches(0.5))

add_text(s4, 'The Insiders Factory changes the economics of character-led content.',
         Inches(0.3), Inches(3.85), Inches(6.0), Inches(0.45),
         font_name='Calibri', size=14, color=DIM)

# Left — comparison table
rows = [
    ('',            'TRADITIONAL',  'INSIDERS FACTORY'),
    ('Time',        '6 – 8 weeks',  '< 10 minutes'),
    ('Cost',        '$12K – $40K',  'Negligible'),
    ('Consistency', 'Variable',     'Guaranteed'),
    ('Volume',      '1 – 2 / month','Unlimited'),
    ('Skill req.',  'Production co.','None'),
]
tx = Inches(0.3)
ty = Inches(4.45)
col_w = [Inches(1.3), Inches(1.8), Inches(2.1)]
row_h = Inches(0.38)

for r, row in enumerate(rows):
    for c, cell in enumerate(row):
        cx = tx + sum(col_w[:c])
        cy = ty + r * row_h
        bg = LIME if r == 0 else (GREY_2 if r % 2 == 0 else GREY_1)
        add_rect(s4, cx, cy, col_w[c], row_h, fill=bg)
        fc = BLACK if r == 0 else (LIME if c == 2 and r > 0 else (WHITE if c == 0 else DIM))
        add_text(s4, cell, cx + Inches(0.1), cy + Inches(0.04),
                 col_w[c] - Inches(0.1), row_h,
                 font_name='Courier New', size=9, bold=(r==0 or c==0),
                 color=fc)

# Right — closing statement
add_rect(s4, Inches(7.0), Inches(1.1), Inches(5.9), Inches(6.1), fill=GREY_1)
add_rect(s4, Inches(7.0), Inches(1.1), Inches(5.9), Emu(12000), fill=LIME)

add_text(s4, 'What this unlocks',
         Inches(7.2), Inches(1.15), Inches(5.5), Inches(0.45),
         font_name='Courier New', size=10, bold=True, color=BLACK)

bullets = [
    '→  A living, breathing brand character that shows up every single week',
    '→  Campaign bursts that can spin up overnight, not over months',
    '→  Consistent character voice and visual identity at any volume',
    '→  Brand and agency teams empowered to move at platform speed',
    '→  A reusable content engine — build it once, run it indefinitely',
]
for i, b in enumerate(bullets):
    add_text(s4, b,
             Inches(7.2), Inches(1.85 + i * 0.82), Inches(5.5), Inches(0.7),
             font_name='Calibri', size=13, color=WHITE if i % 2 == 0 else DIM)

add_text(s4, 'The Insiders Factory',
         Inches(7.2), Inches(6.35), Inches(4.0), Inches(0.5),
         font_name='Georgia', size=20, italic=True, color=LIME)
add_text(s4, 'Built by adam&eveTBWA for IAMS Proactive Health',
         Inches(7.2), Inches(6.78), Inches(5.5), Inches(0.3),
         font_name='Courier New', size=8, color=GREY_3)

add_rect(s4, 0, Inches(7.18), W, Emu(40000), fill=GREY_2)
add_text(s4, 'IAMS PROACTIVE HEALTH  ·  CONFIDENTIAL',
         Inches(0.3), Inches(7.2), Inches(6), Inches(0.28),
         font_name='Courier New', size=7, color=GREY_3)

# ── Save ───────────────────────────────────────────────────────────────────────
out = '/home/user/insiders/insiders-factory-exec-summary.pptx'
prs.save(out)
print(f'Saved → {out}')
